import argparse
import json
import sys
import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def extract_text_from_shape(shape):
    text = ""
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            text += paragraph.text + "\n"
    
    if shape.has_table:
        table = shape.table
        for row in table.rows:
            row_text = []
            for cell in row.cells:
                if cell.text_frame:
                    row_text.append(cell.text_frame.text.strip())
            text += " | ".join(row_text) + "\n"
            
    return text.strip()

def parse_ppt(file_path):
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"Error opening file: {e}", file=sys.stderr)
        sys.exit(1)

    slides_content = []

    for i, slide in enumerate(prs.slides):
        slide_number = i + 1
        slide_data = {
            "slide_number": slide_number,
            "title": "",
            "content": [],
            "notes": ""
        }

        # Extract title
        if slide.shapes.title:
            slide_data["title"] = slide.shapes.title.text.strip()

        # Extract content from shapes
        for shape in slide.shapes:
            # Skip title shape as it's already handled
            if shape == slide.shapes.title:
                continue
            
            text = extract_text_from_shape(shape)
            if text:
                slide_data["content"].append(text)

            # Handle groups recursively if needed (basic implementation here)
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for subshape in shape.shapes:
                    sub_text = extract_text_from_shape(subshape)
                    if sub_text:
                        slide_data["content"].append(sub_text)

        # Extract notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            slide_data["notes"] = slide.notes_slide.notes_text_frame.text.strip()

        slides_content.append(slide_data)

    return slides_content

def format_as_markdown(slides_content):
    md_output = "# Presentation Content\n\n"
    for slide in slides_content:
        md_output += f"## Slide {slide['slide_number']}"
        if slide['title']:
            md_output += f": {slide['title']}"
        md_output += "\n\n"
        
        if slide['content']:
            md_output += "### Content\n"
            for item in slide['content']:
                # Simple bullet point for each text block
                md_output += f"- {item.replace(chr(10), ' ')}\n" # Replace newlines in items to keep list structure clean
            md_output += "\n"
            
        if slide['notes']:
            md_output += "### Notes\n"
            md_output += f"{slide['notes']}\n\n"
            
        md_output += "---\n\n"
    return md_output

def main():
    parser = argparse.ArgumentParser(description="Extract text from PowerPoint files.")
    parser.add_argument("file_path", help="Path to the .pptx file")
    parser.add_argument("--format", choices=["json", "markdown"], default="markdown", help="Output format")
    
    args = parser.parse_args()
    
    if not os.path.exists(args.file_path):
        print(f"File not found: {args.file_path}", file=sys.stderr)
        sys.exit(1)

    content = parse_ppt(args.file_path)
    
    if args.format == "json":
        print(json.dumps(content, indent=2, ensure_ascii=False))
    else:
        print(format_as_markdown(content))

if __name__ == "__main__":
    main()
